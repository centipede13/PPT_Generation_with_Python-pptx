'''
The file takes in images from a folder and makes ppts from image file
'''
import glob
import os
from pptx import Presentation
from pptx.util import Inches

# Input directory for watermarked images
inDir = 'Asset/Watermarked_Images'

# Output directory for resulting ppts
outDir = 'Asset/Images/'

# Function to add images to the slide


def _add_image(slides, image_path):
    left = Inches(1)
    top = Inches(2.5)
    height = Inches(5.5)
    width = Inches(4.5)
    pic = slides.shapes
    pic.add_picture(image_path, left, top, height=height, width=width)


# Main Function
if __name__ == "__main__":
    if os.path.exists(inDir):

        # Presentation Initialized
        prs = Presentation()

        # Slides Height,width
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # Setting Layout
        # For this particular cas I am using Layout 1
        # It has a title Box And content Box
        layout1 = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout1)

        # Looping throgh the folder for jpeg files only
        for jpgFile in glob.glob(os.path.join(inDir + '**/*.jpg')):

            # Setting Title Text,font...
            title = slide.shapes.title
            title.text = 'Sample Title 1'
            title.text_frame.paragraphs[0].font.name = 'Times New Roman'

            # Setting SubTitle Text,font...
            subtitle = slide.placeholders[1]
            subtitle.text = "Sample Subtitle 1"
            subtitle.text_frame.paragraphs[0].font.name = 'Times New Roman'

            # Image add function
            _add_image(slide, jpgFile)

            # extracting only filename example -> image2.jpg from the file
            fileName = os.path.basename(jpgFile)

            # Saving the created ppt with the Imagefile name concatenated with TexT
            prs.save(outDir + '/' + fileName.split('.jpg')
                     [0]+'_MyPresentation.pptx')
